#!/usr/bin/env python3
"""inspect_template.py — Mide un template .pptx para RECONSTRUIRLO fiel en pptxgenjs.

Uso:
    python scripts/inspect_template.py <template.pptx> [outDir]

En estos templates el DISEÑO vive en los slide LAYOUTS (fondos full-bleed como
imágenes, formas de color, y texto); los slides sólo referencian un layout. Por
eso volcamos LAYOUTS y SLIDES, con TODO lo necesario para reproducir el diseño:
geometría, FONDOS, imágenes, y por forma su relleno/borde/tipo + texto enriquecido.

Produce (por defecto en outDir = ".", normalmente la raíz del skill):
    template_spec.json          descripción COMPLETA (pulgadas) por layout y slide
    assets/backgrounds/*        el fondo REAL full-bleed de cada layout (byte-idéntico)
    assets/media/*              logos y decoraciones (deduplicadas por contenido)

`scripts/replica_deck.js` lee este spec y reproduce cada forma tal cual
(`renderFromSpec`). Guía: references/template-replication.md.

Requiere:  pip install python-pptx   (ver scripts/setup.sh)
"""
import hashlib
import json
import os
import sys

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
except Exception as e:  # pragma: no cover
    sys.stderr.write("Falta python-pptx. Instala con: pip install python-pptx\n(%s)\n" % e)
    sys.exit(1)

EMU = 914400.0


def inches(v):
    return round(v / EMU, 3) if v is not None else None


def enum_name(v):
    return str(v).split(".")[-1] if v is not None else None


def as_pt(v):
    """Length → pt; número (multiplicador de interlínea) → float; None → None."""
    if v is None:
        return None
    if hasattr(v, "pt"):
        try:
            return round(float(v.pt), 2)
        except Exception:
            return None
    try:
        return float(v)
    except Exception:
        return None


def color_of(colorformat):
    """RGB explícito, color de tema, o None (heredado)."""
    try:
        if colorformat is None or colorformat.type is None:
            return None
        if getattr(colorformat, "rgb", None) is not None:
            return str(colorformat.rgb)
        if getattr(colorformat, "theme_color", None) is not None:
            return "theme:" + enum_name(colorformat.theme_color)
    except Exception:
        return None
    return None


class ImageSaver:
    """Guarda blobs de imagen deduplicando por contenido; devuelve la ruta relativa."""

    def __init__(self, out):
        self.out = out
        self.by_hash = {}
        os.makedirs(os.path.join(out, "assets", "backgrounds"), exist_ok=True)
        os.makedirs(os.path.join(out, "assets", "media"), exist_ok=True)

    def save(self, blob, ext, subdir, name):
        h = hashlib.md5(blob).hexdigest()
        if h in self.by_hash:
            return self.by_hash[h]
        ext = (ext or "png").lstrip(".")
        fn = "%s.%s" % (name, ext)
        rel = "assets/%s/%s" % (subdir, fn)
        with open(os.path.join(self.out, "assets", subdir, fn), "wb") as fh:
            fh.write(blob)
        self.by_hash[h] = rel
        return rel


def rich_text(shape):
    """Texto completo: párrafos (alineación/spacing) y runs (fuente/color)."""
    if not getattr(shape, "has_text_frame", False):
        return None
    tf = shape.text_frame
    paras = []
    for p in tf.paragraphs:
        runs = []
        for r in p.runs:
            f = r.font
            runs.append(
                {
                    "text": r.text,
                    "face": f.name,
                    "sizePt": float(f.size.pt) if f.size is not None else None,
                    "bold": f.bold,
                    "italic": f.italic,
                    "color": color_of(getattr(f, "color", None)),
                }
            )
        paras.append(
            {
                "align": enum_name(p.alignment),
                "level": p.level,
                "lineSpacing": as_pt(p.line_spacing),
                "spaceBefore": as_pt(p.space_before),
                "spaceAfter": as_pt(p.space_after),
                "runs": runs,
            }
        )
    return {
        "paragraphs": paras,
        "valign": enum_name(getattr(tf, "vertical_anchor", None)),
        "wordWrap": getattr(tf, "word_wrap", None),
    }


def fill_transparency(shape):
    """% de transparencia del relleno solido (0=opaco, 100=invisible), o None.
    Critico: un rectangulo overlay full-bleed suele ir semi-transparente sobre la
    foto; sin esta alfa se reproduce opaco y tapa el fondo."""
    try:
        spPr = shape._element.find(qn("p:spPr"))
        if spPr is None:
            return None
        sf = spPr.find(qn("a:solidFill"))
        if sf is None:
            return None
        clr = sf.find(qn("a:srgbClr"))
        if clr is None:
            clr = sf.find(qn("a:schemeClr"))
        if clr is None:
            return None
        a = clr.find(qn("a:alpha"))
        if a is None:
            return None
        opaque = int(a.get("val")) / 1000.0  # val en milesimas de %
        return round(100.0 - opaque, 1)
    except Exception:
        return None


def shape_fill(shape):
    try:
        f = shape.fill
        t = f.type
    except Exception:
        return None
    if t is None:
        return None
    tn = enum_name(t) or ""
    if "SOLID" in tn:
        try:
            return {"type": "solid", "color": color_of(f.fore_color), "transparency": fill_transparency(shape)}
        except Exception:
            return {"type": "solid", "color": None, "transparency": fill_transparency(shape)}
    if "GRAD" in tn:
        stops = []
        try:
            for gs in f.gradient_stops:
                stops.append(color_of(getattr(gs, "color", None)))
        except Exception:
            pass
        return {"type": "gradient", "stops": stops}
    if "BACKGROUND" in tn:
        return {"type": "none"}
    return {"type": tn.lower()}


def shape_line(shape):
    try:
        ln = shape.line
    except Exception:
        return None
    color = None
    try:
        color = color_of(ln.color)
    except Exception:
        color = None
    width = None
    try:
        width = round(float(ln.width.pt), 2) if ln.width is not None else None
    except Exception:
        width = None
    if color is None and width is None:
        return None
    return {"color": color, "widthPt": width}


def shape_preset(shape):
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return enum_name(shape.auto_shape_type)
    except Exception:
        return None
    return None


def apply_tf(tf, left, top, w, h):
    """Mapea geometria (EMU) de un espacio de coordenadas al absoluto. tf=None → identidad."""
    if left is None or top is None or w is None or h is None or tf is None:
        return left, top, w, h
    gL, gT, sx, sy, chx, chy = tf
    return (gL + (left - chx) * sx, gT + (top - chy) * sy, w * sx, h * sy)


def child_tf_of(grp, g_abs):
    """Transform para los HIJOS de un grupo, dadas sus coords absolutas g_abs (EMU)."""
    try:
        xfrm = grp._element.find(qn("p:grpSpPr") + "/" + qn("a:xfrm"))
        chOff = xfrm.find(qn("a:chOff"))
        chExt = xfrm.find(qn("a:chExt"))
        if chExt is None or None in g_abs:
            return None
        chx, chy = int(chOff.get("x")), int(chOff.get("y"))
        cex, cey = int(chExt.get("cx")), int(chExt.get("cy"))
        gL, gT, gW, gH = g_abs
        sx = gW / cex if cex else 1.0
        sy = gH / cey if cey else 1.0
        return (gL, gT, sx, sy, chx, chy)
    except Exception:
        return None


def walk(shape, tf, out):
    """Aplana grupos: recorre y devuelve (shape, geomAbsEMU) para cada forma-hoja."""
    absg = apply_tf(tf, shape.left, shape.top, shape.width, shape.height)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        ctf = child_tf_of(shape, absg)
        for c in shape.shapes:
            walk(c, ctf, out)
    else:
        out.append((shape, absg))


def shape_record(shape, saver, tag, sw, sh_h, geom):
    kind = enum_name(getattr(shape, "shape_type", None))
    ph = None
    try:
        if shape.is_placeholder:
            ph = enum_name(shape.placeholder_format.type)
    except Exception:
        ph = None
    text = ""
    try:
        text = shape.text if shape.has_text_frame else ""
    except Exception:
        text = ""
    rot = None
    try:
        rot = round(float(shape.rotation), 2) if shape.rotation else None
    except Exception:
        rot = None
    L, T, W, H = geom
    xIn, yIn, wIn, hIn = inches(L), inches(T), inches(W), inches(H)
    rec = {
        "name": getattr(shape, "name", None),
        "kind": kind,
        "placeholder": ph,
        "preset": shape_preset(shape),
        "rotation": rot,
        "xIn": xIn,
        "yIn": yIn,
        "wIn": wIn,
        "hIn": hIn,
        "fill": shape_fill(shape),
        "line": shape_line(shape),
        "text": text,
        "richText": rich_text(shape),
    }
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = shape.image
            w, h, x, y = wIn or 0, hIn or 0, xIn or 0, yIn or 0
            fullbleed = w >= sw * 0.9 and h >= sh_h * 0.9 and x <= 0.2 and y <= 0.2
            sub = "backgrounds" if fullbleed else "media"
            rec["image"] = saver.save(
                img.blob, img.ext, sub, "%s_pic" % tag if fullbleed else "%s_%s" % (tag, shape.shape_id)
            )
            rec["fullbleed"] = fullbleed
        except Exception as e:
            rec["image"] = None
            rec["imageError"] = str(e)
    return rec


def dump_container(container, saver, tag, sw, sh_h):
    leaves = []
    for shp in container.shapes:
        try:
            walk(shp, None, leaves)
        except Exception:
            leaves.append((shp, (shp.left, shp.top, shp.width, shp.height)))
    shapes = []
    background = None
    for shp, absg in leaves:
        try:
            r = shape_record(shp, saver, tag, sw, sh_h, absg)
        except Exception as e:
            r = {"error": str(e), "name": getattr(shp, "name", None)}
        if r.get("fullbleed") and r.get("image") and background is None:
            background = r["image"]
        shapes.append(r)
    return background, shapes


def theme_info(prs):
    from lxml import etree

    major = minor = None
    colors = {}
    try:
        THEME = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        tp = prs.slide_masters[0].part.part_related_by(THEME)
        troot = getattr(tp, "_element", None)
        if troot is None:
            troot = etree.fromstring(tp.blob)
        fs = troot.find(".//" + qn("a:fontScheme"))
        if fs is not None:
            mj = fs.find(qn("a:majorFont") + "/" + qn("a:latin"))
            mn = fs.find(qn("a:minorFont") + "/" + qn("a:latin"))
            major = mj.get("typeface") if mj is not None else None
            minor = mn.get("typeface") if mn is not None else None
        cs = troot.find(".//" + qn("a:clrScheme"))
        if cs is not None:
            for child in cs:
                nm = child.tag.split("}")[-1]
                srgb = child.find(qn("a:srgbClr"))
                sysc = child.find(qn("a:sysClr"))
                if srgb is not None:
                    colors[nm] = srgb.get("val")
                elif sysc is not None:
                    colors[nm] = sysc.get("lastClr") or sysc.get("val")
    except Exception:
        pass
    return {"majorFont": major, "minorFont": minor, "colors": colors}


def main():
    if len(sys.argv) < 2:
        sys.stderr.write(__doc__)
        sys.exit(1)
    template = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) > 2 else "."
    os.makedirs(out, exist_ok=True)
    saver = ImageSaver(out)

    prs = Presentation(template)
    sw = inches(prs.slide_width)
    sh_h = inches(prs.slide_height)

    spec = {
        "source": os.path.basename(template),
        "slideWidthIn": sw,
        "slideHeightIn": sh_h,
        "theme": theme_info(prs),
        "layouts": [],
        "slides": [],
    }

    for i, L in enumerate(prs.slide_layouts, start=1):
        bg, shapes = dump_container(L, saver, "L%d" % i, sw, sh_h)
        spec["layouts"].append({"index": i, "name": getattr(L, "name", None), "background": bg, "shapes": shapes})

    for i, s in enumerate(prs.slides, start=1):
        bg, shapes = dump_container(s, saver, "S%d" % i, sw, sh_h)
        spec["slides"].append(
            {"index": i, "layoutName": getattr(s.slide_layout, "name", None), "background": bg, "shapes": shapes}
        )

    with open(os.path.join(out, "template_spec.json"), "w") as fh:
        json.dump(spec, fh, ensure_ascii=False, indent=2)

    # Preview deck para el ANÁLISIS VISUAL: un slide por layout (add_slide hereda
    # el diseño del layout), para renderizar a thumbnails y verlos con `view`.
    # El diseño de estos templates vive en los layouts, así que renderizar el
    # .pptx original no basta. (Ver references/template-replication.md.)
    n_preview = 0
    try:
        prev = Presentation(template)
        for layout in list(prev.slide_layouts):
            prev.slides.add_slide(layout)
            n_preview += 1
        prev.save(os.path.join(out, "template_layouts_preview.pptx"))
    except Exception as e:
        sys.stderr.write("preview deck omitido: %s\n" % e)
        n_preview = 0

    nbg = len(os.listdir(os.path.join(out, "assets", "backgrounds")))
    nmedia = len(os.listdir(os.path.join(out, "assets", "media")))
    nshapes = sum(len(L["shapes"]) for L in spec["layouts"]) + sum(len(s["shapes"]) for s in spec["slides"])
    print("template_spec.json · %d layouts · %d slides · %d shapes" % (len(spec["layouts"]), len(spec["slides"]), nshapes))
    print("fondos:", nbg, "· media:", nmedia, "· fuentes:", spec["theme"].get("majorFont"), "/", spec["theme"].get("minorFont"))
    if n_preview:
        print("template_layouts_preview.pptx · %d layouts para render/QA visual (ver template-replication.md)" % n_preview)


if __name__ == "__main__":
    main()
